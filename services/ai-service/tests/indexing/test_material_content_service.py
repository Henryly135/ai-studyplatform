from __future__ import annotations

from io import BytesIO
from types import SimpleNamespace
from zipfile import ZipFile

import openpyxl
import pytest
from docx import Document
from pptx import Presentation

from app.services.indexing.material_content_service import MaterialContentRequest, MaterialContentService


def test_extract_text_reads_local_text_file(tmp_path) -> None:
    # Tests local markdown files are read and trimmed into extracted content.
    material_path = tmp_path / "lesson.md"
    material_path.write_text("  # Lesson\n\ncontent  ", encoding="utf-8")

    result = MaterialContentService().extract_text(
        request=MaterialContentRequest(
            title="Lesson",
            content_type="text/markdown",
            storage_provider="LOCAL",
            absolute_path=str(material_path),
            storage_bucket=None,
            object_key="lesson.md",
        )
    )

    assert result.storage_provider == "local"
    assert result.content_text == "# Lesson\n\ncontent"


def test_extract_text_rejects_unsupported_provider() -> None:
    # Tests unsupported material storage providers are rejected.
    with pytest.raises(Exception) as exc_info:
        MaterialContentService().extract_text(
            request=MaterialContentRequest(
                title="Lesson",
                content_type="text/plain",
                storage_provider="s3",
                absolute_path=None,
                storage_bucket=None,
                object_key="lesson.txt",
            )
        )

    assert "Unsupported storageProvider 's3'" in str(exc_info.value)


def test_read_minio_bytes_closes_response(monkeypatch) -> None:
    # Tests MinIO response objects are closed and released after reading.
    response = SimpleNamespace(
        closed=False,
        released=False,
        read=lambda: b"content",
        close=lambda: setattr(response, "closed", True),
        release_conn=lambda: setattr(response, "released", True),
    )
    client = SimpleNamespace(get_object=lambda bucket, key: response)
    monkeypatch.setattr("app.services.indexing.material_content_service.build_minio_client", lambda **_: client)

    payload = MaterialContentService()._read_minio_bytes(bucket_name="bucket", object_key="lesson.txt")

    assert payload == b"content"
    assert response.closed is True
    assert response.released is True


def test_decode_text_falls_back_to_latin1() -> None:
    # Tests text decoding falls back to latin-1 when UTF-8 fails.
    assert MaterialContentService()._decode_text(b"caf\xe9") == "caf\xe9"


def test_extract_text_normalizes_cross_platform_line_endings(tmp_path) -> None:
    material_path = tmp_path / "lesson.txt"
    material_path.write_bytes(b"first\r\nsecond\rthird")

    result = MaterialContentService().extract_text(
        request=MaterialContentRequest(
            title="Lesson",
            content_type="text/plain",
            storage_provider="local",
            absolute_path=str(material_path),
            storage_bucket=None,
            object_key="lesson.txt",
        )
    )

    assert result.content_text == "first\nsecond\nthird"


def test_extract_text_supports_office_and_archive_materials(tmp_path) -> None:
    doc = Document()
    doc.add_paragraph("Word lesson")
    docx_payload = BytesIO()
    doc.save(docx_payload)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Slide lesson"
    pptx_payload = BytesIO()
    presentation.save(pptx_payload)

    workbook = openpyxl.Workbook()
    workbook.active.append(["Sheet lesson", 42])
    xlsx_payload = BytesIO()
    workbook.save(xlsx_payload)

    archive_payload = BytesIO()
    with ZipFile(archive_payload, "w") as archive:
        archive.writestr("lesson.txt", "Archive lesson")

    fixtures = {
        "lesson.docx": (
            docx_payload.getvalue(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "Word lesson",
        ),
        "lesson.pptx": (
            pptx_payload.getvalue(),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "Slide lesson",
        ),
        "lesson.xlsx": (
            xlsx_payload.getvalue(),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "Sheet lesson",
        ),
        "lesson.zip": (archive_payload.getvalue(), "application/zip", "Archive lesson"),
    }

    for filename, (payload, content_type, expected_text) in fixtures.items():
        path = tmp_path / filename
        path.write_bytes(payload)
        result = MaterialContentService().extract_text(
            request=MaterialContentRequest(
                title=filename,
                content_type=content_type,
                storage_provider="local",
                absolute_path=str(path),
                storage_bucket=None,
                object_key=filename,
            )
        )
        assert expected_text in result.content_text


def test_extract_text_routes_bounded_media_to_ai_analysis(monkeypatch, tmp_path) -> None:
    calls: list[dict[str, object]] = []

    def fake_analyze(self, **kwargs):
        calls.append(kwargs)
        return "AI media summary"

    monkeypatch.setattr(
        "app.services.indexing.material_media_analysis_service.MaterialMediaAnalysisService.analyze",
        fake_analyze,
    )
    path = tmp_path / "lesson.png"
    path.write_bytes(b"bounded media")

    result = MaterialContentService().extract_text(
        request=MaterialContentRequest(
            title="Visual lesson",
            content_type="image/png",
            storage_provider="local",
            absolute_path=str(path),
            storage_bucket=None,
            object_key="lesson.png",
        )
    )

    assert result.content_text == "AI media summary"
    assert calls[0]["content_type"] == "image/png"
    assert calls[0]["filename"] == "lesson.png"


def test_extract_text_supports_svg_text_without_remote_upload(tmp_path) -> None:
    path = tmp_path / "diagram.svg"
    path.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg"><title>Architecture</title><text>RAG</text></svg>',
        encoding="utf-8",
    )

    result = MaterialContentService().extract_text(
        request=MaterialContentRequest(
            title="Diagram",
            content_type="image/svg+xml",
            storage_provider="local",
            absolute_path=str(path),
            storage_bucket=None,
            object_key=path.name,
        )
    )

    assert "Architecture" in result.content_text
    assert "RAG" in result.content_text
