from __future__ import annotations

import re
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass
from html.parser import HTMLParser
from io import BytesIO
from pathlib import Path
from xml.etree import ElementTree

import openpyxl
import xlrd
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy.orm import Session

from app.core.config import settings
from app.services.indexing.material_media_analysis_service import MaterialMediaAnalysisService
from platform_common.errors import invalid_request_error
from platform_common.storage import build_minio_client

_TEXT_CONTENT_TYPES = {
    "text/plain",
    "text/markdown",
    "application/json",
    "text/csv",
    "text/tab-separated-values",
    "text/yaml",
    "application/yaml",
    "application/x-yaml",
    "text/html",
    "application/xhtml+xml",
    "text/xml",
    "application/xml",
    "text/rtf",
    "application/rtf",
}
_WORD_CONTENT_TYPES = {
    "application/msword",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}
_PRESENTATION_CONTENT_TYPES = {
    "application/vnd.ms-powerpoint",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
_SPREADSHEET_CONTENT_TYPES = {
    "application/vnd.ms-excel",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}
_ODF_CONTENT_TYPES = {
    "application/vnd.oasis.opendocument.text",
    "application/vnd.oasis.opendocument.presentation",
    "application/vnd.oasis.opendocument.spreadsheet",
}
_MEDIA_PREFIXES = ("image/", "audio/", "video/")


class _VisibleTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        normalized = " ".join(data.split())
        if normalized:
            self.parts.append(normalized)


@dataclass(frozen=True)
class MaterialContentRequest:
    title: str
    content_type: str | None
    storage_provider: str
    absolute_path: str | None
    storage_bucket: str | None
    object_key: str


@dataclass(frozen=True)
class ExtractedMaterialContent:
    title: str
    content_type: str | None
    storage_provider: str
    object_key: str
    content_text: str


class MaterialContentService:
    def __init__(self, session: Session | None = None) -> None:
        self.session = session
        self.media_analysis = MaterialMediaAnalysisService(session)

    def extract_text(self, *, request: MaterialContentRequest) -> ExtractedMaterialContent:
        normalized_provider = request.storage_provider.strip().lower()
        payload = self._read_bytes(request=request, storage_provider=normalized_provider)
        normalized_content_type = (request.content_type or "").strip().lower() or None
        content_text = self._extract_text_payload(
            payload=payload,
            content_type=normalized_content_type,
            object_key=request.object_key,
            title=request.title,
        )
        content_text = self._limit_text(content_text.replace("\r\n", "\n").replace("\r", "\n"))
        if not content_text.strip():
            raise invalid_request_error("Extracted material content is empty")

        return ExtractedMaterialContent(
            title=request.title,
            content_type=request.content_type,
            storage_provider=normalized_provider,
            object_key=request.object_key,
            content_text=content_text.strip(),
        )

    def _read_bytes(self, *, request: MaterialContentRequest, storage_provider: str) -> bytes:
        if storage_provider == "local":
            return self._read_local_bytes(absolute_path=request.absolute_path)
        if storage_provider == "minio":
            return self._read_minio_bytes(
                bucket_name=request.storage_bucket,
                object_key=request.object_key,
            )
        raise invalid_request_error(f"Unsupported storageProvider '{request.storage_provider}'")

    def _read_local_bytes(self, *, absolute_path: str | None) -> bytes:
        if not absolute_path or not absolute_path.strip():
            raise invalid_request_error("absolutePath is required for local material indexing")
        target_path = Path(absolute_path)
        if not target_path.is_file():
            raise invalid_request_error(f"Material file not found at '{absolute_path}'")
        return target_path.read_bytes()

    def _read_minio_bytes(self, *, bucket_name: str | None, object_key: str) -> bytes:
        if not bucket_name or not bucket_name.strip():
            raise invalid_request_error("storageBucket is required for MinIO material indexing")

        client = build_minio_client(
            endpoint=settings.minio_endpoint,
            access_key=settings.minio_access_key,
            secret_key=settings.minio_secret_key,
        )
        response = client.get_object(bucket_name.strip(), object_key)
        try:
            return response.read()
        finally:
            response.close()
            response.release_conn()

    def _extract_text_payload(
        self,
        *,
        payload: bytes,
        content_type: str | None,
        object_key: str,
        title: str,
        archive_depth: int = 0,
    ) -> str:
        suffix = Path(object_key).suffix.lower()
        if content_type == "application/pdf" or suffix == ".pdf":
            return self._extract_pdf_text(payload)
        # Legacy binary Office formats need the command-line converters.  Check
        # the extension before the broad MIME groups because browsers report
        # .doc/.ppt/.xls with the same MIME families as their OOXML variants.
        if suffix == ".doc":
            return self._extract_legacy_office_text(payload=payload, suffix=suffix)
        if suffix == ".ppt":
            return self._extract_legacy_office_text(payload=payload, suffix=suffix)
        if suffix == ".xls":
            return self._extract_xls_text(payload)
        if content_type in _WORD_CONTENT_TYPES or suffix == ".docx":
            return self._extract_docx_text(payload)
        if content_type in _PRESENTATION_CONTENT_TYPES or suffix == ".pptx":
            return self._extract_pptx_text(payload)
        if content_type in _SPREADSHEET_CONTENT_TYPES or suffix in {".xlsx", ".xlsm"}:
            return self._extract_xlsx_text(payload)
        if content_type in _ODF_CONTENT_TYPES or suffix in {".odt", ".odp", ".ods"}:
            return self._extract_odf_text(payload)
        if suffix == ".svg":
            return self._extract_svg_text(payload=payload, title=title)
        if content_type in {"application/zip", "application/x-zip-compressed"} or suffix == ".zip":
            if archive_depth >= 2:
                return ""
            return self._extract_archive_text(payload, archive_depth=archive_depth)
        if content_type and content_type.startswith(_MEDIA_PREFIXES):
            return self.media_analysis.analyze(
                title=title,
                content_type=content_type,
                filename=Path(object_key).name,
                payload=payload,
            )
        if suffix in {".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v", ".mpeg", ".mpg"}:
            return self.media_analysis.analyze(
                title=title,
                content_type=content_type or "video/mp4",
                filename=Path(object_key).name,
                payload=payload,
            )
        if suffix in {".mp3", ".wav", ".aac", ".m4a", ".ogg", ".flac", ".opus", ".wma"}:
            return self.media_analysis.analyze(
                title=title,
                content_type=content_type or "audio/mpeg",
                filename=Path(object_key).name,
                payload=payload,
            )
        if suffix in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff", ".svg"}:
            return self.media_analysis.analyze(
                title=title,
                content_type=content_type or "image/png",
                filename=Path(object_key).name,
                payload=payload,
            )
        if content_type in _TEXT_CONTENT_TYPES or suffix in {
            ".txt",
            ".md",
            ".markdown",
            ".csv",
            ".tsv",
            ".json",
            ".yaml",
            ".yml",
            ".xml",
            ".html",
            ".htm",
        }:
            return self._decode_text(payload, content_type=content_type, suffix=suffix)
        raise invalid_request_error(
            f"Unsupported material content type '{content_type or suffix or 'unknown'}' for text extraction"
        )

    def _extract_pdf_text(self, payload: bytes) -> str:
        reader = PdfReader(BytesIO(payload))
        page_text: list[str] = []
        for page in reader.pages:
            page_text.append((page.extract_text() or "").strip())
        return "\n\n".join(text for text in page_text if text)

    def _extract_docx_text(self, payload: bytes) -> str:
        document = Document(BytesIO(payload))
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        table_rows = [" | ".join(cell.text.strip() for cell in row.cells) for table in document.tables for row in table.rows]
        return "\n\n".join([*paragraphs, *[row for row in table_rows if row]])

    def _extract_pptx_text(self, payload: bytes) -> str:
        presentation = Presentation(BytesIO(payload))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    value = str(shape.text).strip()
                    if value:
                        texts.append(value)
            if texts:
                slides.append(f"第 {index} 页\n" + "\n".join(texts))
        return "\n\n".join(slides)

    def _extract_xlsx_text(self, payload: bytes) -> str:
        workbook = openpyxl.load_workbook(BytesIO(payload), read_only=True, data_only=True)
        try:
            sheets: list[str] = []
            for sheet in workbook.worksheets:
                rows: list[str] = []
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                    if values:
                        rows.append(" | ".join(values))
                    if len(rows) >= 5000:
                        break
                if rows:
                    sheets.append(f"工作表：{sheet.title}\n" + "\n".join(rows))
            return "\n\n".join(sheets)
        finally:
            workbook.close()

    def _extract_xls_text(self, payload: bytes) -> str:
        workbook = xlrd.open_workbook(file_contents=payload, on_demand=True)
        sheets: list[str] = []
        try:
            for sheet in workbook.sheets():
                rows: list[str] = []
                for row_index in range(min(sheet.nrows, 5000)):
                    values = [str(value).strip() for value in sheet.row_values(row_index) if str(value).strip()]
                    if values:
                        rows.append(" | ".join(values))
                if rows:
                    sheets.append(f"工作表：{sheet.name}\n" + "\n".join(rows))
            return "\n\n".join(sheets)
        finally:
            workbook.release_resources()

    def _extract_odf_text(self, payload: bytes) -> str:
        try:
            with zipfile.ZipFile(BytesIO(payload)) as archive:
                content = archive.read("content.xml")
        except (KeyError, zipfile.BadZipFile) as exc:
            raise invalid_request_error("Unable to read OpenDocument content") from exc
        try:
            root = ElementTree.fromstring(content)
        except ElementTree.ParseError as exc:
            raise invalid_request_error("Unable to parse OpenDocument content") from exc
        return " ".join(part.strip() for part in root.itertext() if part and part.strip())

    def _extract_svg_text(self, *, payload: bytes, title: str) -> str:
        try:
            root = ElementTree.fromstring(payload)
        except ElementTree.ParseError as exc:
            raise invalid_request_error("Unable to parse SVG material") from exc
        visible_text = " ".join(part.strip() for part in root.itertext() if part and part.strip())
        return visible_text or f"SVG 图片资料：{title}"

    def _extract_legacy_office_text(self, *, payload: bytes, suffix: str) -> str:
        command = "antiword" if suffix == ".doc" else "catppt"
        temporary_file = tempfile.NamedTemporaryFile(prefix="material-office-", suffix=suffix, delete=False)
        path = Path(temporary_file.name)
        try:
            temporary_file.write(payload)
            temporary_file.flush()
            temporary_file.close()
            completed = subprocess.run(
                [command, str(path)],
                capture_output=True,
                text=True,
                check=False,
                timeout=60,
            )
        except (OSError, subprocess.TimeoutExpired) as exc:
            raise invalid_request_error(f"{suffix} 文件解析工具不可用") from exc
        finally:
            path.unlink(missing_ok=True)
        if completed.returncode != 0:
            raise invalid_request_error(f"无法解析 {suffix} 文件")
        return completed.stdout

    def _extract_archive_text(self, payload: bytes, *, archive_depth: int) -> str:
        try:
            archive = zipfile.ZipFile(BytesIO(payload))
        except zipfile.BadZipFile as exc:
            raise invalid_request_error("Unable to read ZIP material") from exc
        extracted: list[str] = []
        total_uncompressed = 0
        try:
            members = [info for info in archive.infolist() if not info.is_dir()]
            if len(members) > settings.ai_material_archive_max_entries:
                raise invalid_request_error("ZIP material contains too many files")
            for info in members:
                total_uncompressed += info.file_size
                if total_uncompressed > settings.ai_material_archive_max_uncompressed_bytes:
                    raise invalid_request_error("ZIP material expands beyond the configured analysis limit")
                name = Path(info.filename)
                if name.is_absolute() or ".." in name.parts:
                    continue
                suffix = name.suffix.lower()
                if suffix in {".mp4", ".mov", ".avi", ".mkv", ".webm", ".mp3", ".wav", ".aac", ".m4a", ".ogg", ".flac"}:
                    continue
                try:
                    member_payload = archive.read(info)
                    member_text = self._extract_text_payload(
                        payload=member_payload,
                        content_type=None,
                        object_key=str(name),
                        title=name.name,
                        archive_depth=archive_depth + 1,
                    )
                except Exception:
                    continue
                if member_text.strip():
                    extracted.append(f"文件：{name.as_posix()}\n{member_text.strip()}")
            return "\n\n".join(extracted)
        finally:
            archive.close()

    def _decode_text(
        self,
        payload: bytes,
        *,
        content_type: str | None = None,
        suffix: str = "",
    ) -> str:
        for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
            try:
                decoded = payload.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            raise invalid_request_error("Unable to decode material text content")
        if content_type in {"text/html", "application/xhtml+xml"} or suffix in {".html", ".htm", ".xhtml"}:
            parser = _VisibleTextParser()
            parser.feed(decoded)
            return " ".join(parser.parts)
        if content_type in {"text/rtf", "application/rtf"} or suffix == ".rtf":
            return re.sub(r"\\[a-zA-Z]+-?\d* ?|[{}]", "", decoded)
        return decoded

    def _limit_text(self, content_text: str) -> str:
        normalized = content_text.strip()
        if len(normalized) <= settings.ai_material_max_extracted_chars:
            return normalized
        return normalized[: settings.ai_material_max_extracted_chars] + "\n[资料内容已按索引上限截断]"
